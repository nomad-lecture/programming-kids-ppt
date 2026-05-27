import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(44, 62, 80)
BLUE = RGBColor(52, 152, 219)
GREEN = RGBColor(39, 174, 96)
ORANGE = RGBColor(243, 156, 18)
RED = RGBColor(231, 76, 60)
DARK = RGBColor(33, 33, 33)
LIGHT_BG = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)

IMG_DIR = "images"


def img_path(name):
    return os.path.join(IMG_DIR, name)


def add_bg(slide, color=LIGHT_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, text="Програмчлалын танилцуулга"):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.0), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.color.rgb = BLUE


def add_note(slide, note_text):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = note_text


def add_title_box(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.2), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = NAVY

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.05), Inches(11.8), Inches(0.5))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(15)
        r2.font.color.rgb = DARK


def add_bullets(slide, bullets, left=0.8, top=1.7, width=6.0, height=4.8, font_size=22):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        for r in p.runs:
            r.font.size = Pt(font_size)
            r.font.color.rgb = DARK


def add_image_or_placeholder(slide, image_file, left, top, width, height, label="Зураг"):
    full_path = img_path(image_file)
    if os.path.exists(full_path):
        slide.shapes.add_picture(full_path, Inches(left), Inches(top), Inches(width), Inches(height))
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = BLUE
        shape.line.width = Pt(2)

        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"{label}\n({image_file} not found)"
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = BLUE


def add_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(232, 244, 252))

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.0), Inches(11.7), Inches(4.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BLUE
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.clear()

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = "Програмчлал гэж юу вэ?"
    r1.font.size = Pt(30)
    r1.font.bold = True
    r1.font.color.rgb = NAVY

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "5-р ангийн хүүхдүүдэд зориулсан танилцуулга"
    r2.font.size = Pt(18)
    r2.font.color.rgb = DARK

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = "[Таны нэр] — Программист"
    r3.font.size = Pt(18)
    r3.font.color.rgb = GREEN

    p4 = tf.add_paragraph()
    p4.alignment = PP_ALIGN.CENTER
    r4 = p4.add_run()
    r4.text = "Roblox • Minecraft • Код • Тоглоом"
    r4.font.size = Pt(16)
    r4.font.color.rgb = ORANGE

    add_footer(slide)
    add_note(slide, "Өөрийгөө танилцуул. Өнөөдрийн сэдвээ хэл.")


def add_two_col_slide(title, bullets, image_file, label, note_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_box(slide, title)
    add_bullets(slide, bullets, left=0.8, top=1.7, width=6.1, height=4.8, font_size=22)
    add_image_or_placeholder(slide, image_file, left=7.4, top=1.8, width=5.0, height=3.8, label=label)
    add_footer(slide)
    add_note(slide, note_text)


def add_simple_slide(title, bullets, note_text, accent=BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.9))
    banner.fill.solid()
    banner.fill.fore_color.rgb = accent
    banner.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = WHITE

    add_bullets(slide, bullets, left=0.9, top=1.4, width=11.2, height=5.4, font_size=23)
    add_footer(slide)
    add_note(slide, note_text)


add_title_slide()

add_two_col_slide(
    "Та нарын хэн Roblox, Minecraft тоглодог вэ?",
    [
        "Гар өргөөрэй",
        "Эдгээр тоглоомуудын ард юу байдаг бол?",
        "Хэн ийм тоглоомыг хийдэг вэ?"
    ],
    "roblox_minecraft.png",
    "Roblox / Minecraft зураг",
    "Хүүхдүүдийг оролцуулж эхэл."
)

add_simple_slide(
    "Програмчлал гэж юу вэ?",
    [
        "Програмчлал гэдэг нь компьютерт заавар өгөхийг хэлнэ",
        "Алхам алхмаар ойлгомжтой хэлж өгнө",
        "Компьютер яг өгсөн зааврыг дагадаг"
    ],
    "Энгийн тайлбар өг."
)

add_simple_slide(
    "Алгоритм гэж юу вэ?",
    [
        "Алгоритм = алхам алхмын заавар",
        "Жишээ: өглөө босоод сургууль руу явах",
        "Дараалал зөв байх хэрэгтэй"
    ],
    "Алгоритмыг өдөр тутмын жишээгээр тайлбарла.",
    ORANGE
)

add_two_col_slide(
    "Би программист хүн",
    [
        "Компьютерт заавар өгдөг",
        "Шинэ зүйл бүтээдэг",
        "Алдаа олж засдаг",
        "Хүмүүст хэрэгтэй програм дээр ажилладаг"
    ],
    "code.png",
    "Кодны зураг",
    "Өөрийн ажлыг хүүхдийн хэлээр тайлбарла."
)

add_two_col_slide(
    "Roblox, Minecraft-ийн ард код байдаг",
    [
        "Дүр хөдөлнө",
        "Үсэрнэ",
        "Хаалга онгойно",
        "Оноо нэмэгдэнэ",
        "Даалгавар эхэлнэ"
    ],
    "minecraft.png",
    "Minecraft зураг",
    "Тоглоомын ард код байдгийг ярь."
)

add_two_col_slide(
    "Код дандаа урт, хэцүү байх албагүй",
    [
        "Зарим код block хэлбэртэй байдаг",
        "Scratch шиг хэрэгслээр хүүхдүүд ч сурч болно",
        "Код бол заавар өгөх арга юм"
    ],
    "scratch.png",
    "Scratch / block coding зураг",
    "Block coding-ийн тухай товч хэл."
)

add_simple_slide(
    "Үндсэн ойлголтууд",
    [
        "1. Дараалал",
        "2. Давталт",
        "3. Нөхцөл",
        "4. Алдаа засах"
    ],
    "4 ойлголтоо тани��цуул.",
    RED
)

add_simple_slide(
    "Тоглоом 1: Багшийг программчилъя",
    [
        "Надад яг таг заавар өгөөрэй",
        "Урагш 2 алх",
        "Зүүн эргэ",
        "Номоо ав",
        "Ширээн дээр тавь"
    ],
    "Робот болж тогло.",
    GREEN
)

add_simple_slide(
    "Тоглоом 2: Хэрвээ ... бол ...",
    [
        "Хэрвээ би 'нар' гэвэл бос",
        "Хэрвээ би 'бороо' гэвэл суу",
        "Хэрвээ би 'салхи' гэвэл нэг эргэ",
        "Хэрвээ би 'цас' гэвэл толгойгоо дар"
    ],
    "Condition тоглоом тоглуул.",
    BLUE
)

add_simple_slide(
    "Тоглоом 3: Алдаа олоорой",
    [
        "Би буруу дараалал хэлнэ",
        "Та нар алдааг нь олно",
        "Дараа нь зөв болгож засна"
    ],
    "Debugging ойлголтыг тоглоомоор ойлгуул.",
    ORANGE
)

add_simple_slide(
    "Дүгнэлт",
    [
        "Програмчлал = компьютерт заавар өгөх",
        "Код = шинэ зүйл бүтээх",
        "Алдаа = сурах боломж",
        "Та нар ч бас сурч чадна"
    ],
    "Эцсийн урам өг.",
    RED
)

add_simple_slide(
    "Асуулт байна уу?",
    [
        "Та нар юу асуумаар байна?"
    ],
    "Асуулт хариулт ав."
)

prs.save("programming_for_kids_with_images.pptx")
print("Created: programming_for_kids_with_images.pptx")
